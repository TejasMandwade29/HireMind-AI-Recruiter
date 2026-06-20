import os
import shutil
import win32com.client
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- PATHS SETUP ---
WORKSPACE_DIR = r"c:\Users\Dell\Downloads\[PUB] India_runs_data_and_ai_challenge\[PUB] India_runs_data_and_ai_challenge\India_runs_data_and_ai_challenge"
ARTIFACTS_DIR = r"C:\Users\Dell\.gemini\antigravity-ide\brain\45c24265-7a76-4c4c-a037-19d50c2804c1"
TEMPLATE_PPTX = os.path.join(WORKSPACE_DIR, "ppt_requierment", "Idea Submission Template _ Redrob.pptx")
OUTPUT_PPTX = os.path.join(WORKSPACE_DIR, "hiremind_presentation.pptx")

SCREENSHOT_SRC = os.path.join(WORKSPACE_DIR, "ppt_requierment", "Screenshot 2026-06-20 114118.png")
VISUALIZER_DST = os.path.join(WORKSPACE_DIR, "score_visualizer.png")
VERDICT_DST = os.path.join(WORKSPACE_DIR, "recruiter_verdict.png")

# --- IMAGE CROPPING STEP ---
print("Cropping screenshot assets programmatically...")
if os.path.exists(SCREENSHOT_SRC):
    img = Image.open(SCREENSHOT_SRC)
    # Crop right half for score contribution bars (910x320)
    cropped_vis = img.crop((950, 260, 1860, 580))
    cropped_vis.save(VISUALIZER_DST)
    
    # Crop left half for AI recruiter verdict text (910x320)
    cropped_verd = img.crop((20, 260, 930, 580))
    cropped_verd.save(VERDICT_DST)
    print("Symmetric crops generated successfully.")
else:
    print(f"Error: Screenshot source {SCREENSHOT_SRC} does not exist!")

# --- STYLE CONSTANTS (HireMind Dark Palette for Inner Cards) ---
BG_COLOR = RGBColor(15, 23, 42)        # Deep Navy (#0F172A)
CARD_BG_COLOR = RGBColor(15, 23, 42)    # Deep Navy (#0F172A)
TEXT_WHITE = RGBColor(255, 255, 255)   # White
TEXT_CYAN = RGBColor(0, 242, 254)     # Electric Cyan (#00F2FE)
TEXT_SLATE = RGBColor(148, 163, 184)   # Slate Gray (#94A3B8)
TEXT_CRIMSON = RGBColor(239, 68, 68)   # Crimson (#EF4444)
CARD_BORDER_COLOR = RGBColor(71, 85, 105) # Slate Border (#475569)

print(f"Loading official template: {TEMPLATE_PPTX}...")
prs = Presentation(TEMPLATE_PPTX)

def add_metric_badge(slide, badge_text):
    # Add a pill shape for consistent metric badge in top right (x=6.8, y=0.89, w=2.7, h=0.45)
    badge_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(0.89), Inches(2.7), Inches(0.4)
    )
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = CARD_BG_COLOR
    badge_bg.line.color.rgb = TEXT_CYAN
    badge_bg.line.width = Pt(1.0)
    
    tf = badge_bg.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = badge_text.upper()
    p.font.name = 'Manrope'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_CYAN

# ==============================================================================
# SLIDE 2: Solution Overview (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 2...")
slide_2 = prs.slides[1]

# 1. Update Title (Shape 1 text only)
for shape in slide_2.shapes:
    if shape.has_text_frame and "solution overview" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Solution Overview"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders safely
for shape in list(slide_2.shapes):
    if shape.has_text_frame and "proposed solution" in shape.text.lower():
        spTree = slide_2.shapes._spTree
        spTree.remove(shape._element)

# 3. Hero Headline Text Box — increased whitespace (Top: 1.15)
hero_box = slide_2.shapes.add_textbox(Inches(0.41), Inches(1.15), Inches(4.0), Inches(0.5))
tf_hero = hero_box.text_frame
tf_hero.word_wrap = True
tf_hero.margin_left = tf_hero.margin_top = tf_hero.margin_right = tf_hero.margin_bottom = 0

p_hero = tf_hero.paragraphs[0]
p_hero.text = "From 100,000 Candidates to the Right Hire"
p_hero.font.name = 'Manrope'
p_hero.font.size = Pt(17)
p_hero.font.bold = True
p_hero.font.color.rgb = RGBColor(15, 23, 42)

# 4. Left Card A: Traditional ATS Column — MUTED GRAY border
ATS_BORDER_COLOR = RGBColor(100, 116, 139)  # Muted slate gray
ATS_TITLE_COLOR = RGBColor(148, 163, 184)   # Slate text

ats_card = slide_2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.41), Inches(1.80), Inches(2.0), Inches(3.0)
)
ats_card.fill.solid()
ats_card.fill.fore_color.rgb = CARD_BG_COLOR
ats_card.line.color.rgb = ATS_BORDER_COLOR
ats_card.line.width = Pt(1.0)

tf_ats = ats_card.text_frame
tf_ats.word_wrap = True
tf_ats.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_ats.margin_left = Inches(0.12)
tf_ats.margin_right = Inches(0.12)
tf_ats.margin_top = Inches(0.12)
tf_ats.margin_bottom = Inches(0.12)

p_ats_h = tf_ats.paragraphs[0]
p_ats_h.alignment = PP_ALIGN.CENTER
p_ats_h.text = "TRADITIONAL ATS"
p_ats_h.font.name = 'Manrope'
p_ats_h.font.size = Pt(10)
p_ats_h.font.bold = True
p_ats_h.font.color.rgb = ATS_TITLE_COLOR
p_ats_h.space_after = Pt(10)

ats_rows = [
    ("Keyword Matching", "Simple text filters"),
    ("Easy to Game", "Susceptible to tricks"),
    ("Fake Profiles", "Profiles slip through"),
    ("Black Box Ranks", "Unexplained logic")
]
for title, desc in ats_rows:
    p_tr = tf_ats.add_paragraph()
    p_tr.text = title
    p_tr.font.name = 'Manrope'
    p_tr.font.size = Pt(9)
    p_tr.font.bold = True
    p_tr.font.color.rgb = TEXT_WHITE
    
    p_dr = tf_ats.add_paragraph()
    p_dr.text = desc
    p_dr.font.name = 'Manrope'
    p_dr.font.size = Pt(7.5)
    p_dr.font.color.rgb = TEXT_SLATE
    p_dr.space_after = Pt(4)

# 5. Left Card B: HireMind Column
hm_card = slide_2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2.51), Inches(1.80), Inches(2.0), Inches(3.0)
)
hm_card.fill.solid()
hm_card.fill.fore_color.rgb = CARD_BG_COLOR
hm_card.line.color.rgb = TEXT_CYAN
hm_card.line.width = Pt(1.5)

tf_hm = hm_card.text_frame
tf_hm.word_wrap = True
tf_hm.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_hm.margin_left = Inches(0.12)
tf_hm.margin_right = Inches(0.12)
tf_hm.margin_top = Inches(0.12)
tf_hm.margin_bottom = Inches(0.12)

p_hm_h = tf_hm.paragraphs[0]
p_hm_h.alignment = PP_ALIGN.CENTER
p_hm_h.text = "HIREMIND ENGINE"
p_hm_h.font.name = 'Manrope'
p_hm_h.font.size = Pt(10)
p_hm_h.font.bold = True
p_hm_h.font.color.rgb = TEXT_CYAN
p_hm_h.space_after = Pt(10)

hm_rows = [
    ("Smart Ranking", "Looks beyond keywords"),
    ("Fraud Checks", "Detects suspicious profiles"),
    ("Anomaly Detection", "Filters spam & traps"),
    ("Clear Rankings", "Shows why candidates rank higher")
]
for title, desc in hm_rows:
    p_tr = tf_hm.add_paragraph()
    p_tr.text = title
    p_tr.font.name = 'Manrope'
    p_tr.font.size = Pt(9)
    p_tr.font.bold = True
    p_tr.font.color.rgb = TEXT_CYAN
    
    p_dr = tf_hm.add_paragraph()
    p_dr.text = desc
    p_dr.font.name = 'Manrope'
    p_dr.font.size = Pt(7.5)
    p_dr.font.color.rgb = TEXT_WHITE
    p_dr.space_after = Pt(4)

# 6. Right Column: Funnel — DOMINANT VISUAL (wider, taller, bigger numbers)
# Badge aligned with funnel
add_metric_badge(slide_2, "100,000 Candidates Analyzed")

funnel_steps = [
    {"num": "100,000", "label": "RAW APPLICANTS", "height": 0.55},
    {"num": "81,101", "label": "CLEAN CANDIDATES (Anomaly-free)", "height": 0.55},
    {"num": "1,092", "label": "TIER 1 & 2 SPECIALISTS MATCHED", "height": 0.55},
]

y_pos = 1.80
for idx, step in enumerate(funnel_steps):
    box_h = step["height"]
    
    node = slide_2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(y_pos), Inches(5.0), Inches(box_h)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = CARD_BORDER_COLOR if idx < 2 else CARD_BG_COLOR
    node.line.fill.background()
    
    tf_node = node.text_frame
    tf_node.word_wrap = True
    tf_node.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_node.margin_left = Inches(0.1)
    tf_node.margin_right = Inches(0.1)
    tf_node.margin_top = tf_node.margin_bottom = 0
    
    p_num = tf_node.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = step["num"]
    p_num.font.name = 'Manrope'
    p_num.font.size = Pt(20)
    p_num.font.bold = True
    p_num.font.color.rgb = TEXT_WHITE
    
    p_lbl = tf_node.add_paragraph()
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.text = step["label"]
    p_lbl.font.name = 'Manrope'
    p_lbl.font.size = Pt(7.5)
    p_lbl.font.color.rgb = TEXT_SLATE
    
    # Down arrow
    arrow_box = slide_2.shapes.add_textbox(Inches(7.05), Inches(y_pos + box_h), Inches(0.4), Inches(0.18))
    tf_arrow = arrow_box.text_frame
    tf_arrow.margin_left = tf_arrow.margin_top = tf_arrow.margin_right = tf_arrow.margin_bottom = 0
    p_arr = tf_arrow.paragraphs[0]
    p_arr.alignment = PP_ALIGN.CENTER
    p_arr.text = "\u25bc"
    p_arr.font.size = Pt(8)
    p_arr.font.color.rgb = TEXT_SLATE
    
    y_pos += box_h + 0.22

# 7. Final TOP 100 Card — Dark navy bg + cyan text/border (premium, not solid cyan)
top100_y = y_pos
top100_card = slide_2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(top100_y), Inches(5.0), Inches(0.85)
)
top100_card.fill.solid()
top100_card.fill.fore_color.rgb = CARD_BG_COLOR
top100_card.line.color.rgb = TEXT_CYAN
top100_card.line.width = Pt(2.0)

tf_top = top100_card.text_frame
tf_top.word_wrap = True
tf_top.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_top.margin_left = Inches(0.1)
tf_top.margin_right = Inches(0.1)
tf_top.margin_top = tf_top.margin_bottom = 0

p_top_num = tf_top.paragraphs[0]
p_top_num.alignment = PP_ALIGN.CENTER
r_trophy = p_top_num.add_run()
r_trophy.text = "\U0001F3C6 "
r_trophy.font.size = Pt(22)
r_num = p_top_num.add_run()
r_num.text = "TOP 100"
r_num.font.name = 'Manrope'
r_num.font.size = Pt(24)
r_num.font.bold = True
r_num.font.color.rgb = TEXT_CYAN

p_top_lbl = tf_top.add_paragraph()
p_top_lbl.alignment = PP_ALIGN.CENTER
p_top_lbl.text = "RECRUITER-READY SHORTLIST"
p_top_lbl.font.name = 'Manrope'
p_top_lbl.font.size = Pt(9)
p_top_lbl.font.bold = True
p_top_lbl.font.color.rgb = TEXT_WHITE

# ==============================================================================
# SLIDE 3: Beyond Keyword Matching (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 3...")
slide_3 = prs.slides[2]

for shape in slide_3.shapes:
    if shape.has_text_frame and "jd understanding" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Beyond Keyword Matching"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

for shape in list(slide_3.shapes):
    if shape.has_text_frame and "extracted from" in shape.text.lower():
        spTree = slide_3.shapes._spTree
        spTree.remove(shape._element)

add_metric_badge(slide_3, "5 Recruiter Signals Evaluated")

sub_box = slide_3.shapes.add_textbox(Inches(0.41), Inches(1.40), Inches(9.32), Inches(0.4))
tf_sub = sub_box.text_frame
tf_sub.word_wrap = True
tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0

p_sub = tf_sub.paragraphs[0]
p_sub.text = "Candidates are ranked using five hiring signals."
p_sub.font.name = 'Manrope'
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = RGBColor(15, 23, 42)

# Left Column Stacked Proof Images
vis_border = slide_3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.40), Inches(1.84), Inches(4.32), Inches(1.22)
)
vis_border.fill.solid()
vis_border.fill.fore_color.rgb = CARD_BG_COLOR
vis_border.line.color.rgb = TEXT_CYAN
vis_border.line.width = Pt(1.5)

if os.path.exists(VISUALIZER_DST):
    slide_3.shapes.add_picture(VISUALIZER_DST, Inches(0.41), Inches(1.85), Inches(4.3), Inches(1.2))

verd_border = slide_3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.40), Inches(3.19), Inches(4.32), Inches(1.22)
)
verd_border.fill.solid()
verd_border.fill.fore_color.rgb = CARD_BG_COLOR
verd_border.line.color.rgb = TEXT_CYAN
verd_border.line.width = Pt(1.5)

if os.path.exists(VERDICT_DST):
    slide_3.shapes.add_picture(VERDICT_DST, Inches(0.41), Inches(3.20), Inches(4.3), Inches(1.2))

caption_box = slide_3.shapes.add_textbox(Inches(0.41), Inches(4.50), Inches(4.3), Inches(0.25))
tf_cap = caption_box.text_frame
tf_cap.word_wrap = True
tf_cap.margin_left = tf_cap.margin_top = tf_cap.margin_right = tf_cap.margin_bottom = 0
p_cap = tf_cap.paragraphs[0]
p_cap.text = "Actual HireMind output showing transparent candidate evaluation and ranking explanation."
p_cap.font.name = 'Manrope'
p_cap.font.size = Pt(7.5)
p_cap.font.italic = True
p_cap.font.color.rgb = TEXT_SLATE

badge_border = slide_3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.41), Inches(4.80), Inches(1.6), Inches(0.25)
)
badge_border.fill.solid()
badge_border.fill.fore_color.rgb = CARD_BG_COLOR
badge_border.line.color.rgb = TEXT_CYAN
badge_border.line.width = Pt(1.0)

tf_bd = badge_border.text_frame
tf_bd.margin_left = tf_bd.margin_top = tf_bd.margin_right = tf_bd.margin_bottom = 0
p_bd = tf_bd.paragraphs[0]
p_bd.alignment = PP_ALIGN.CENTER
p_bd.text = "✓ Transparent Rankings"
p_bd.font.name = 'Manrope'
p_bd.font.size = Pt(8.0)
p_bd.font.bold = True
p_bd.font.color.rgb = TEXT_CYAN

# Right Column Signal Cards
signal_cards = [
    {
        "title": "Template Fit", "weight": "25%",
        "desc": "Matches candidates to the role, not just keywords."
    },
    {
        "title": "Semantic Skills", "weight": "25%",
        "desc": "Measures how closely skills match the job."
    },
    {
        "title": "Candidate Intent", "weight": "20%",
        "desc": "Measures candidate availability and interest."
    },
    {
        "title": "Growth", "weight": "15%",
        "desc": "Tracks career growth over time."
    },
    {
        "title": "Stability", "weight": "15%",
        "desc": "Rewards stable work history."
    }
]

for idx, card in enumerate(signal_cards):
    y_pos = 1.85 + (idx * 0.59)
    
    card_shape = slide_3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.9), Inches(y_pos), Inches(4.8), Inches(0.54)
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = CARD_BG_COLOR
    card_shape.line.color.rgb = CARD_BORDER_COLOR
    card_shape.line.width = Pt(1.0)
    
    card_text = slide_3.shapes.add_textbox(Inches(5.0), Inches(y_pos + 0.05), Inches(4.6), Inches(0.44))
    tf_card = card_text.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = tf_card.margin_top = tf_card.margin_right = tf_card.margin_bottom = 0
    
    p_header = tf_card.paragraphs[0]
    
    run_weight = p_header.add_run()
    run_weight.text = f"{card['weight']}   "
    run_weight.font.name = 'Manrope'
    run_weight.font.size = Pt(14)
    run_weight.font.bold = True
    run_weight.font.color.rgb = TEXT_CYAN
    
    run_title = p_header.add_run()
    run_title.text = card['title']
    run_title.font.name = 'Manrope'
    run_title.font.size = Pt(11)
    run_title.font.bold = True
    run_title.font.color.rgb = TEXT_WHITE
    
    p_desc = tf_card.add_paragraph()
    p_desc.text = card["desc"]
    p_desc.font.name = 'Manrope'
    p_desc.font.size = Pt(8.5)
    p_desc.font.color.rgb = TEXT_SLATE
    p_desc.space_before = Pt(1)

# ==============================================================================
# SLIDE 4: Ranking Methodology (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 4...")
slide_4 = prs.slides[3]

# 1. Update Title (Shape 1 text only)
for shape in slide_4.shapes:
    if shape.has_text_frame and "ranking methodology" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Ranking Methodology"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders safely
for shape in list(slide_4.shapes):
    if shape.has_text_frame and "retrieve, score" in shape.text.lower():
        spTree = slide_4.shapes._spTree
        spTree.remove(shape._element)

# 3. Add Metric Badge
add_metric_badge(slide_4, "100% Local Processing")

# 4. Hero Title Text Box (Left: 0.41, Top: 1.10, Width: 9.32, Height: 0.3)
hero4_box = slide_4.shapes.add_textbox(Inches(0.41), Inches(1.10), Inches(9.32), Inches(0.3))
tf_hero4 = hero4_box.text_frame
tf_hero4.word_wrap = True
tf_hero4.margin_left = tf_hero4.margin_top = tf_hero4.margin_right = tf_hero4.margin_bottom = 0
p_h4 = tf_hero4.paragraphs[0]
p_h4.text = "From Job Description to Top 100 Shortlist"
p_h4.font.name = 'Manrope'
p_h4.font.size = Pt(18)
p_h4.font.bold = True
p_h4.font.color.rgb = RGBColor(15, 23, 42)

# Subtitle Text Box (Left: 0.41, Top: 1.45, Width: 9.32, Height: 0.25)
sub4_box = slide_4.shapes.add_textbox(Inches(0.41), Inches(1.45), Inches(9.32), Inches(0.25))
tf_sub4 = sub4_box.text_frame
tf_sub4.word_wrap = True
tf_sub4.margin_left = tf_sub4.margin_top = tf_sub4.margin_right = tf_sub4.margin_bottom = 0
p_sub4 = tf_sub4.paragraphs[0]
p_sub4.text = "Ranks candidates locally in under 70 seconds."
p_sub4.font.name = 'Manrope'
p_sub4.font.size = Pt(12)
p_sub4.font.color.rgb = TEXT_SLATE

# 5. Hero Visual: Horizontal Ranking Pipeline Nodes (Top: 1.85, Height: 1.80)
pipeline_nodes_v2 = [
    {
        "title": "Job Description",
        "desc": "Read the job requirements."
    },
    {
        "title": "Candidate Matching",
        "desc": "Find matching candidates."
    },
    {
        "title": "Multi-Signal Evaluation",
        "desc": "Score each candidate."
    },
    {
        "title": "Final Ranking",
        "desc": "Create the final ranking."
    },
    {
        "title": "Top 100 Shortlist",
        "desc": "Top candidates ready for review."
    }
]

for idx, step in enumerate(pipeline_nodes_v2):
    x_pos = 0.41 + (idx * 1.88)
    is_last = (idx == 4)
    
    # Node Background shape
    node_color = TEXT_CYAN if is_last else CARD_BG_COLOR
    node_text_color = BG_COLOR if is_last else TEXT_WHITE
    node_desc_color = BG_COLOR if is_last else TEXT_SLATE
    
    node = slide_4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.85), Inches(1.6), Inches(1.8)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = node_color
    node.line.color.rgb = TEXT_CYAN if is_last else CARD_BORDER_COLOR
    node.line.width = Pt(1.5 if is_last else 1.0)
    
    tf_node = node.text_frame
    tf_node.word_wrap = True
    tf_node.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_node.margin_left = Inches(0.1)
    tf_node.margin_right = Inches(0.1)
    tf_node.margin_top = Inches(0.1)
    tf_node.margin_bottom = Inches(0.1)
    
    # Title
    p_n = tf_node.paragraphs[0]
    p_n.alignment = PP_ALIGN.CENTER
    p_n.text = step["title"]
    p_n.font.name = 'Manrope'
    p_n.font.size = Pt(11)
    p_n.font.bold = True
    p_n.font.color.rgb = node_text_color
    p_n.space_after = Pt(8)
    
    # Description
    p_d = tf_node.add_paragraph()
    p_d.alignment = PP_ALIGN.CENTER
    p_d.text = step["desc"]
    p_d.font.name = 'Manrope'
    p_d.font.size = Pt(8.5)
    p_d.font.color.rgb = node_desc_color
    
    # Arrow between nodes
    if idx < 4:
        arrow_box = slide_4.shapes.add_textbox(Inches(x_pos + 1.6), Inches(2.55), Inches(0.28), Inches(0.4))
        tf_arr = arrow_box.text_frame
        tf_arr.margin_left = tf_arr.margin_right = tf_arr.margin_top = tf_arr.margin_bottom = 0
        p_a = tf_arr.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        p_a.text = "➔"
        p_a.font.name = 'Manrope'
        p_a.font.size = Pt(16)
        p_a.font.color.rgb = TEXT_SLATE

# 6. Bottom Cards: Constraint validation & Model details (Top: 3.95, Height: 1.0)
# Left Compliance Card
const_bg = slide_4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.41), Inches(3.95), Inches(4.3), Inches(1.0)
)
const_bg.fill.solid()
const_bg.fill.fore_color.rgb = CARD_BG_COLOR
const_bg.line.color.rgb = TEXT_CYAN
const_bg.line.width = Pt(1.5)

tf_co = const_bg.text_frame
tf_co.word_wrap = True
tf_co.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_co.margin_left = Inches(0.2)
tf_co.margin_right = Inches(0.2)
tf_co.margin_top = Inches(0.1)
tf_co.margin_bottom = Inches(0.1)

# Clear default paragraph
tf_co.paragraphs[0].text = ""

items = ["CPU Only", "No External APIs", "<70s Runtime"]
for idx, item in enumerate(items):
    p = tf_co.add_paragraph() if idx > 0 else tf_co.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    
    # Checkmark run
    r_check = p.add_run()
    r_check.text = "✓  "
    r_check.font.name = 'Manrope'
    r_check.font.size = Pt(10)
    r_check.font.bold = True
    r_check.font.color.rgb = TEXT_CYAN
    
    # Text run
    r_text = p.add_run()
    r_text.text = item
    r_text.font.name = 'Manrope'
    r_text.font.size = Pt(10)
    r_text.font.bold = True
    r_text.font.color.rgb = TEXT_WHITE

# Right Model Card
model_bg = slide_4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.29), Inches(3.95), Inches(4.3), Inches(1.0)
)
model_bg.fill.solid()
model_bg.fill.fore_color.rgb = CARD_BG_COLOR
model_bg.line.color.rgb = CARD_BORDER_COLOR
model_bg.line.width = Pt(1.0)

tf_mo = model_bg.text_frame
tf_mo.word_wrap = True
tf_mo.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_mo.margin_left = Inches(0.2)
tf_mo.margin_right = Inches(0.2)
tf_mo.margin_top = Inches(0.1)
tf_mo.margin_bottom = Inches(0.1)

tf_mo.paragraphs[0].text = ""

# Title
p_t = tf_mo.paragraphs[0]
p_t.alignment = PP_ALIGN.CENTER
p_t.space_after = Pt(4)
r_title = p_t.add_run()
r_title.text = "LOCAL EMBEDDING MODEL"
r_title.font.name = 'Manrope'
r_title.font.size = Pt(8.5)
r_title.font.bold = True
r_title.font.color.rgb = TEXT_SLATE

# Value
p_v = tf_mo.add_paragraph()
p_v.alignment = PP_ALIGN.CENTER
r_val = p_v.add_run()
r_val.text = "all-MiniLM-L6-v2"
r_val.font.name = 'Manrope'
r_val.font.size = Pt(13)
r_val.font.bold = True
r_val.font.color.rgb = TEXT_CYAN

# ==============================================================================
# SLIDE 5: Explainability & Data Validation (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 5...")
slide_5 = prs.slides[4]

# 1. Update Title (Shape 1 text only)
for shape in slide_5.shapes:
    if shape.has_text_frame and "explainability" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Explainability & Data Validation"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders safely
for shape in list(slide_5.shapes):
    if shape.has_text_frame and "prevent hallucinations" in shape.text.lower():
        spTree = slide_5.shapes._spTree
        spTree.remove(shape._element)

# 3. Add Metric Badge
add_metric_badge(slide_5, "Fact-Based Explanations")

# 4. Hero Subtitle Text Box (Trustworthy Rankings)
trust_box = slide_5.shapes.add_textbox(Inches(0.41), Inches(1.40), Inches(9.32), Inches(0.4))
tf_trust = trust_box.text_frame
tf_trust.word_wrap = True
tf_trust.margin_left = tf_trust.margin_top = tf_trust.margin_right = tf_trust.margin_bottom = 0

p_trust = tf_trust.paragraphs[0]
p_trust.text = "Trustworthy & Explainable Rankings"
p_trust.font.name = 'Manrope'
p_trust.font.size = Pt(16)
p_trust.font.bold = True
p_trust.font.color.rgb = RGBColor(15, 23, 42)

# 5. Left Column: Data Integrity (Left: 0.41, Width: 4.3, Top: 1.85, Height: 2.9)
integ_card = slide_5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.41), Inches(1.85), Inches(4.3), Inches(2.9)
)
integ_card.fill.solid()
integ_card.fill.fore_color.rgb = CARD_BG_COLOR
integ_card.line.color.rgb = TEXT_CYAN
integ_card.line.width = Pt(1.5)

tf_in = integ_card.text_frame
tf_in.word_wrap = True
tf_in.margin_left = Inches(0.2)
tf_in.margin_top = Inches(0.18)

p_in_h = tf_in.paragraphs[0]
p_in_h.text = "18,899 INVALID PROFILES FILTERED"
p_in_h.font.name = 'Manrope'
p_in_h.font.size = Pt(16)
p_in_h.font.bold = True
p_in_h.font.color.rgb = TEXT_CYAN
p_in_h.space_after = Pt(4)

p_in_sub = tf_in.add_paragraph()
p_in_sub.text = "Suspicious and low-quality profiles removed before ranking."
p_in_sub.font.name = 'Manrope'
p_in_sub.font.size = Pt(9.5)
p_in_sub.font.color.rgb = TEXT_SLATE
p_in_sub.space_after = Pt(20)

bullets_in = [
    "✔  Experience Validation",
    "✔  Career History Checks",
    "✔  Honeypot Detection"
]
for bullet in bullets_in:
    p_b = tf_in.add_paragraph()
    p_b.text = bullet
    p_b.font.name = 'Manrope'
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = TEXT_WHITE
    p_b.space_after = Pt(10)

p_in_foot = tf_in.add_paragraph()
p_in_foot.text = "(Checks profile data for inconsistencies.)"
p_in_foot.font.name = 'Manrope'
p_in_foot.font.size = Pt(8.5)
p_in_foot.font.italic = True
p_in_foot.font.color.rgb = TEXT_SLATE
p_in_foot.space_before = Pt(8)

# 6. Right Column: AI Recruiter Verdict Screenshot (Left: 5.0, Top: 1.85, Width: 4.73, Height: 1.66)
img5_border = slide_5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.85), Inches(4.73), Inches(1.66)
)
img5_border.fill.solid()
img5_border.fill.fore_color.rgb = CARD_BG_COLOR
img5_border.line.color.rgb = TEXT_CYAN
img5_border.line.width = Pt(1.5)

if os.path.exists(VERDICT_DST):
    print("Inserting verdict screenshot on Slide 5...")
    slide_5.shapes.add_picture(
        VERDICT_DST, Inches(5.01), Inches(1.86), Inches(4.71), Inches(1.64)
    )
else:
    print("Warning: recruiter_verdict.png not found, skipping picture insert.")

# Right Column Bottom Subtext Caption (Left: 5.0, Top: 3.65, Width: 4.73, Height: 1.1)
cap5_box = slide_5.shapes.add_textbox(Inches(5.0), Inches(3.60), Inches(4.73), Inches(1.15))
tf_cap5 = cap5_box.text_frame
tf_cap5.word_wrap = True
tf_cap5.margin_left = tf_cap5.margin_top = tf_cap5.margin_right = tf_cap5.margin_bottom = 0

p_cap5_h = tf_cap5.paragraphs[0]
p_cap5_h.text = "Factual AI Reasoning"
p_cap5_h.font.name = 'Manrope'
p_cap5_h.font.size = Pt(12)
p_cap5_h.font.bold = True
p_cap5_h.font.color.rgb = TEXT_CYAN
p_cap5_h.space_after = Pt(2)

p_cap5_b = tf_cap5.add_paragraph()
p_cap5_b.text = "Fact-based reasoning generated directly from candidate data."
p_cap5_b.font.name = 'Manrope'
p_cap5_b.font.size = Pt(10.5)
p_cap5_b.font.bold = True
p_cap5_b.font.color.rgb = TEXT_WHITE
p_cap5_b.space_after = Pt(4)

p_cap5_d = tf_cap5.add_paragraph()
p_cap5_d.text = "Every recommendation is based on candidate data, score signals, and profile history."
p_cap5_d.font.name = 'Manrope'
p_cap5_d.font.size = Pt(8.5)
p_cap5_d.font.color.rgb = TEXT_SLATE
p_cap5_d.space_after = Pt(8)

p_cap5_bd = tf_cap5.add_paragraph()
p_cap5_bd.text = "✓ Transparent Rankings   |   ✓ Explainable Decisions"
p_cap5_bd.font.name = 'Manrope'
p_cap5_bd.font.size = Pt(8.5)
p_cap5_bd.font.bold = True
p_cap5_bd.font.color.rgb = TEXT_CYAN

# ==============================================================================
# SLIDE 6: End-to-End Workflow (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 6...")
slide_6 = prs.slides[5]

# 1. Update Title
for shape in slide_6.shapes:
    if shape.has_text_frame and "end-to-end workflow" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "End-to-End Workflow"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders
for shape in list(slide_6.shapes):
    if shape.has_text_frame and "complete workflow" in shape.text.lower():
        spTree = slide_6.shapes._spTree
        spTree.remove(shape._element)

# 3. Add Metric Badge
add_metric_badge(slide_6, "Recruiter-Ready Pipeline")

# 4. Hero Title
hero6_box = slide_6.shapes.add_textbox(Inches(0.41), Inches(1.10), Inches(9.32), Inches(0.3))
tf_hero6 = hero6_box.text_frame
tf_hero6.word_wrap = True
tf_hero6.margin_left = tf_hero6.margin_top = tf_hero6.margin_right = tf_hero6.margin_bottom = 0
p_h6 = tf_hero6.paragraphs[0]
p_h6.text = "From Job Input to Recruiter Shortlist"
p_h6.font.name = 'Manrope'
p_h6.font.size = Pt(18)
p_h6.font.bold = True
p_h6.font.color.rgb = RGBColor(15, 23, 42)

# 5. Three Workflow Stage Cards
workflow_stages = [
    {
        "number": "1",
        "title": "Enter Job Requirements",
        "desc": "Recruiter pastes or enters the job description directly into the secure local dashboard.",
        "highlight": False
    },
    {
        "number": "2",
        "title": "Analyze Candidates",
        "desc": "System filters suspicious profiles, matches skills locally, and computes the 5-signal score.",
        "highlight": False
    },
    {
        "number": "3",
        "title": "Review & Shortlist",
        "desc": "Recruiter inspects transparent ranking explanations and selects top candidates.",
        "highlight": True
    }
]

for idx, stage in enumerate(workflow_stages):
    x_pos = 0.41 + (idx * 3.15)
    is_highlight = stage["highlight"]
    card_color = TEXT_CYAN if is_highlight else CARD_BG_COLOR
    title_color = BG_COLOR if is_highlight else TEXT_WHITE
    desc_color = BG_COLOR if is_highlight else TEXT_SLATE
    num_color = BG_COLOR if is_highlight else TEXT_CYAN
    
    # Card background
    card = slide_6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.60), Inches(2.85), Inches(2.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = card_color
    card.line.color.rgb = TEXT_CYAN if is_highlight else CARD_BORDER_COLOR
    card.line.width = Pt(1.5 if is_highlight else 1.0)
    
    tf_card = card.text_frame
    tf_card.word_wrap = True
    tf_card.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_card.margin_left = Inches(0.2)
    tf_card.margin_right = Inches(0.2)
    tf_card.margin_top = Inches(0.15)
    tf_card.margin_bottom = Inches(0.15)
    
    # Step number
    p_num = tf_card.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = f"STEP {stage['number']}"
    p_num.font.name = 'Manrope'
    p_num.font.size = Pt(10)
    p_num.font.bold = True
    p_num.font.color.rgb = num_color
    p_num.space_after = Pt(8)
    
    # Title
    p_title = tf_card.add_paragraph()
    p_title.alignment = PP_ALIGN.CENTER
    p_title.text = stage["title"]
    p_title.font.name = 'Manrope'
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(10)
    
    # Description
    p_desc = tf_card.add_paragraph()
    p_desc.alignment = PP_ALIGN.CENTER
    p_desc.text = stage["desc"]
    p_desc.font.name = 'Manrope'
    p_desc.font.size = Pt(9)
    p_desc.font.color.rgb = desc_color
    
    # Arrow between cards
    if idx < 2:
        arrow_box = slide_6.shapes.add_textbox(Inches(x_pos + 2.85), Inches(2.7), Inches(0.3), Inches(0.4))
        tf_arr = arrow_box.text_frame
        tf_arr.margin_left = tf_arr.margin_right = tf_arr.margin_top = tf_arr.margin_bottom = 0
        p_a = tf_arr.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        p_a.text = "\u2794"
        p_a.font.name = 'Manrope'
        p_a.font.size = Pt(16)
        p_a.font.color.rgb = TEXT_SLATE

# 6. Small note under Step 2 area
note6_box = slide_6.shapes.add_textbox(Inches(0.41), Inches(4.35), Inches(9.32), Inches(0.25))
tf_note6 = note6_box.text_frame
tf_note6.word_wrap = True
tf_note6.margin_left = tf_note6.margin_top = tf_note6.margin_right = tf_note6.margin_bottom = 0
p_note6 = tf_note6.paragraphs[0]
p_note6.alignment = PP_ALIGN.CENTER
p_note6.text = "Completed in under 70 seconds on standard hardware."
p_note6.font.name = 'Manrope'
p_note6.font.size = Pt(9.5)
p_note6.font.italic = True
p_note6.font.color.rgb = TEXT_SLATE

# ==============================================================================
# SLIDE 7: System Architecture (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 7...")
slide_7 = prs.slides[6]

# 1. Update Title
for shape in slide_7.shapes:
    if shape.has_text_frame and "system architecture" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "System Architecture"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Add Metric Badge
add_metric_badge(slide_7, "Fully Offline System")

# 3. Hero Title
hero7_box = slide_7.shapes.add_textbox(Inches(0.41), Inches(1.10), Inches(9.32), Inches(0.3))
tf_hero7 = hero7_box.text_frame
tf_hero7.word_wrap = True
tf_hero7.margin_left = tf_hero7.margin_top = tf_hero7.margin_right = tf_hero7.margin_bottom = 0
p_h7 = tf_hero7.paragraphs[0]
p_h7.text = "Modular Design Optimized for Offline Performance"
p_h7.font.name = 'Manrope'
p_h7.font.size = Pt(18)
p_h7.font.bold = True
p_h7.font.color.rgb = RGBColor(15, 23, 42)

# 4. Left Column: 3-Layer Architecture Stack (Width: 5.4)
arch_layers = [
    {
        "title": "PRESENTATION LAYER",
        "component": "Streamlit UI",
        "desc": "Secure local control panel and explainable candidate review dashboard."
    },
    {
        "title": "INTELLIGENCE LAYER",
        "component": "Ranking & Validation Engine",
        "desc": "Consistent candidate ranking, multi-signal weights, and fraud filtering."
    },
    {
        "title": "DATA LAYER",
        "component": "Local Data Store",
        "desc": "JSONL profiles, JD keywords, and offline embedding store."
    }
]

for idx, layer in enumerate(arch_layers):
    y_pos = 1.55 + (idx * 1.05)
    
    # Layer card
    layer_card = slide_7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.41), Inches(y_pos), Inches(5.4), Inches(0.90)
    )
    layer_card.fill.solid()
    layer_card.fill.fore_color.rgb = CARD_BG_COLOR
    layer_card.line.color.rgb = TEXT_CYAN if idx == 1 else CARD_BORDER_COLOR
    layer_card.line.width = Pt(1.5 if idx == 1 else 1.0)
    
    tf_layer = layer_card.text_frame
    tf_layer.word_wrap = True
    tf_layer.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_layer.margin_left = Inches(0.2)
    tf_layer.margin_right = Inches(0.15)
    tf_layer.margin_top = Inches(0.08)
    tf_layer.margin_bottom = Inches(0.08)
    
    # Layer title
    p_lt = tf_layer.paragraphs[0]
    p_lt.text = layer["title"]
    p_lt.font.name = 'Manrope'
    p_lt.font.size = Pt(8)
    p_lt.font.bold = True
    p_lt.font.color.rgb = TEXT_SLATE
    p_lt.space_after = Pt(2)
    
    # Component name
    p_comp = tf_layer.add_paragraph()
    p_comp.text = layer["component"]
    p_comp.font.name = 'Manrope'
    p_comp.font.size = Pt(13)
    p_comp.font.bold = True
    p_comp.font.color.rgb = TEXT_CYAN
    p_comp.space_after = Pt(2)
    
    # Description
    p_ld = tf_layer.add_paragraph()
    p_ld.text = layer["desc"]
    p_ld.font.name = 'Manrope'
    p_ld.font.size = Pt(8.5)
    p_ld.font.color.rgb = TEXT_WHITE
    
    # Down arrow between layers
    if idx < 2:
        arrow_y = y_pos + 0.90
        # Data flow labels
        flow_label = "Ranked Results \u2193" if idx == 0 else "Candidate Data \u2191"
        arrow_box = slide_7.shapes.add_textbox(Inches(2.5), Inches(arrow_y), Inches(1.5), Inches(0.15))
        tf_arr = arrow_box.text_frame
        tf_arr.margin_left = tf_arr.margin_right = tf_arr.margin_top = tf_arr.margin_bottom = 0
        p_a = tf_arr.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        p_a.text = flow_label
        p_a.font.name = 'Manrope'
        p_a.font.size = Pt(7.5)
        p_a.font.bold = True
        p_a.font.color.rgb = TEXT_SLATE

# 5. Right Column: Engine Constraints Panel (Width: 3.7)
constraints_card = slide_7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(1.55), Inches(3.6), Inches(3.0)
)
constraints_card.fill.solid()
constraints_card.fill.fore_color.rgb = CARD_BG_COLOR
constraints_card.line.color.rgb = TEXT_CYAN
constraints_card.line.width = Pt(1.5)

tf_con = constraints_card.text_frame
tf_con.word_wrap = True
tf_con.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_con.margin_left = Inches(0.2)
tf_con.margin_right = Inches(0.2)
tf_con.margin_top = Inches(0.15)
tf_con.margin_bottom = Inches(0.15)

# Title
p_con_h = tf_con.paragraphs[0]
p_con_h.alignment = PP_ALIGN.CENTER
p_con_h.text = "STRICT HARDWARE COMPLIANCE"
p_con_h.font.name = 'Manrope'
p_con_h.font.size = Pt(9)
p_con_h.font.bold = True
p_con_h.font.color.rgb = TEXT_WHITE
p_con_h.space_after = Pt(16)

constraint_items = [
    ("100% Offline Execution", "No external APIs or data leakage."),
    ("CPU-Only Optimization", "Runs on consumer laptop hardware."),
    ("Standard Python Environment", "Zero complex cloud dependencies.")
]

for title, desc in constraint_items:
    p_ci = tf_con.add_paragraph()
    p_ci.space_after = Pt(4)
    
    r_check = p_ci.add_run()
    r_check.text = "\u2713  "
    r_check.font.name = 'Manrope'
    r_check.font.size = Pt(10)
    r_check.font.bold = True
    r_check.font.color.rgb = TEXT_CYAN
    
    r_title = p_ci.add_run()
    r_title.text = title
    r_title.font.name = 'Manrope'
    r_title.font.size = Pt(10)
    r_title.font.bold = True
    r_title.font.color.rgb = TEXT_WHITE
    
    p_cd = tf_con.add_paragraph()
    p_cd.text = desc
    p_cd.font.name = 'Manrope'
    p_cd.font.size = Pt(8)
    p_cd.font.color.rgb = TEXT_SLATE
    p_cd.space_after = Pt(10)

# ==============================================================================
# SLIDE 8: Results & Performance (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 8...")
slide_8 = prs.slides[7]

# 1. Update Title
for shape in slide_8.shapes:
    if shape.has_text_frame and shape.text.strip() == "Results & Performance":
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Results & Performance"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders
for shape in list(slide_8.shapes):
    if shape.has_text_frame and "runtime and compute" in shape.text.lower():
        spTree = slide_8.shapes._spTree
        spTree.remove(shape._element)

# 3. Add Metric Badge
add_metric_badge(slide_8, "Challenge Compliant")

# 4. Hero Title
hero8_box = slide_8.shapes.add_textbox(Inches(0.41), Inches(1.10), Inches(9.32), Inches(0.3))
tf_hero8 = hero8_box.text_frame
tf_hero8.word_wrap = True
tf_hero8.margin_left = tf_hero8.margin_top = tf_hero8.margin_right = tf_hero8.margin_bottom = 0
p_h8 = tf_hero8.paragraphs[0]
p_h8.text = "High-Speed Offline Matching Under 70 Seconds"
p_h8.font.name = 'Manrope'
p_h8.font.size = Pt(18)
p_h8.font.bold = True
p_h8.font.color.rgb = RGBColor(15, 23, 42)

# 5. Left Hero Card: Runtime (the strongest proof point)
hero_card = slide_8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.41), Inches(1.55), Inches(4.3), Inches(3.4)
)
hero_card.fill.solid()
hero_card.fill.fore_color.rgb = CARD_BG_COLOR
hero_card.line.color.rgb = TEXT_CYAN
hero_card.line.width = Pt(1.5)

tf_hero_card = hero_card.text_frame
tf_hero_card.word_wrap = True
tf_hero_card.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_hero_card.margin_left = Inches(0.25)
tf_hero_card.margin_right = Inches(0.25)
tf_hero_card.margin_top = Inches(0.15)
tf_hero_card.margin_bottom = Inches(0.15)

# Hero number — largest font on the slide
p_hero_num = tf_hero_card.paragraphs[0]
p_hero_num.alignment = PP_ALIGN.CENTER
p_hero_num.text = "68.2s"
p_hero_num.font.name = 'Manrope'
p_hero_num.font.size = Pt(54)
p_hero_num.font.bold = True
p_hero_num.font.color.rgb = TEXT_CYAN
p_hero_num.space_after = Pt(4)

p_hero_label = tf_hero_card.add_paragraph()
p_hero_label.alignment = PP_ALIGN.CENTER
p_hero_label.text = "End-to-End Processing Time"
p_hero_label.font.name = 'Manrope'
p_hero_label.font.size = Pt(14)
p_hero_label.font.bold = True
p_hero_label.font.color.rgb = TEXT_WHITE
p_hero_label.space_after = Pt(8)

p_hero_note = tf_hero_card.add_paragraph()
p_hero_note.alignment = PP_ALIGN.CENTER
p_hero_note.text = "Meets challenge runtime requirements."
p_hero_note.font.name = 'Manrope'
p_hero_note.font.size = Pt(9.5)
p_hero_note.font.italic = True
p_hero_note.font.color.rgb = TEXT_SLATE

# 6. Right Column: 2x2 Metric Grid
metric_grid = [
    {"value": "100,000", "label": "Candidates Evaluated", "row": 0, "col": 0},
    {"value": "18,899", "label": "Invalid Profiles Removed", "row": 0, "col": 1},
    {"value": "1,092", "label": "High-Quality Matches", "row": 1, "col": 0},
    {"value": "0%", "label": "Honeypots in Top 100", "row": 1, "col": 1}
]

for metric in metric_grid:
    x = 5.0 + (metric["col"] * 2.42)
    y = 1.55 + (metric["row"] * 1.62)
    
    m_card = slide_8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.28), Inches(1.48)
    )
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = CARD_BG_COLOR
    m_card.line.color.rgb = CARD_BORDER_COLOR
    m_card.line.width = Pt(1.0)
    
    tf_m = m_card.text_frame
    tf_m.word_wrap = True
    tf_m.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_m.margin_left = Inches(0.12)
    tf_m.margin_right = Inches(0.12)
    tf_m.margin_top = Inches(0.1)
    tf_m.margin_bottom = Inches(0.1)
    
    p_mv = tf_m.paragraphs[0]
    p_mv.alignment = PP_ALIGN.CENTER
    p_mv.text = metric["value"]
    p_mv.font.name = 'Manrope'
    p_mv.font.size = Pt(24)
    p_mv.font.bold = True
    p_mv.font.color.rgb = TEXT_CYAN
    p_mv.space_after = Pt(4)
    
    p_ml = tf_m.add_paragraph()
    p_ml.alignment = PP_ALIGN.CENTER
    p_ml.text = metric["label"]
    p_ml.font.name = 'Manrope'
    p_ml.font.size = Pt(9)
    p_ml.font.bold = True
    p_ml.font.color.rgb = TEXT_WHITE

# 7. Bottom quality statement
qual_box = slide_8.shapes.add_textbox(Inches(5.0), Inches(4.80), Inches(4.7), Inches(0.2))
tf_qual = qual_box.text_frame
tf_qual.word_wrap = True
tf_qual.margin_left = tf_qual.margin_top = tf_qual.margin_right = tf_qual.margin_bottom = 0
p_qual = tf_qual.paragraphs[0]
p_qual.alignment = PP_ALIGN.CENTER
p_qual.text = "Every shortlisted candidate includes transparent ranking explanations."
p_qual.font.name = 'Manrope'
p_qual.font.size = Pt(8)
p_qual.font.italic = True
p_qual.font.color.rgb = TEXT_SLATE

# ==============================================================================
# SLIDE 9: Technologies Used (Edit-in-place)
# ==============================================================================
print("Rebuilding Slide 9...")
slide_9 = prs.slides[8]

# 1. Update Title
for shape in slide_9.shapes:
    if shape.has_text_frame and "technologies used" in shape.text.lower():
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Technologies Used"
        p.font.name = 'Manrope'
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)

# 2. Remove default template placeholders
for shape in list(slide_9.shapes):
    if shape.has_text_frame and "technologies, frameworks" in shape.text.lower():
        spTree = slide_9.shapes._spTree
        spTree.remove(shape._element)

# 3. Add Metric Badge
add_metric_badge(slide_9, "Zero Cloud Dependencies")

# 4. Hero Title
hero9_box = slide_9.shapes.add_textbox(Inches(0.41), Inches(1.10), Inches(9.32), Inches(0.3))
tf_hero9 = hero9_box.text_frame
tf_hero9.word_wrap = True
tf_hero9.margin_left = tf_hero9.margin_top = tf_hero9.margin_right = tf_hero9.margin_bottom = 0
p_h9 = tf_hero9.paragraphs[0]
p_h9.text = "Lightweight Stack Selected for Compliance and Speed"
p_h9.font.name = 'Manrope'
p_h9.font.size = Pt(18)
p_h9.font.bold = True
p_h9.font.color.rgb = RGBColor(15, 23, 42)

# 5. Technology Category Cards (2x2 grid)
tech_cards = [
    {
        "category": "INTERFACE",
        "techs": "Streamlit",
        "reason": "Chosen for rapid local deployment and clean interactive review.",
        "row": 0, "col": 0
    },
    {
        "category": "PROCESSING",
        "techs": "Python  \u2022  Pandas  \u2022  NumPy",
        "reason": "Selected for high-speed, in-memory filtering of 100k candidate rows.",
        "row": 0, "col": 1
    },
    {
        "category": "SKILL MATCHING",
        "techs": "Sentence Transformers\nall-MiniLM-L6-v2",
        "reason": "Selected for high-accuracy, offline skill matching on CPU.",
        "row": 1, "col": 0
    },
    {
        "category": "CUSTOM RANKING ENGINE",
        "techs": "Multi-signal scoring\nValidation rules\nExplainability generation",
        "reason": "Built for consistent, multi-signal evaluation with zero API dependencies.",
        "row": 1, "col": 1
    }
]

for tc in tech_cards:
    x = 0.41 + (tc["col"] * 4.75)
    y = 1.55 + (tc["row"] * 1.62)
    
    t_card = slide_9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(4.5), Inches(1.48)
    )
    t_card.fill.solid()
    t_card.fill.fore_color.rgb = CARD_BG_COLOR
    t_card.line.color.rgb = CARD_BORDER_COLOR
    t_card.line.width = Pt(1.0)
    
    tf_tc = t_card.text_frame
    tf_tc.word_wrap = True
    tf_tc.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_tc.margin_left = Inches(0.2)
    tf_tc.margin_right = Inches(0.15)
    tf_tc.margin_top = Inches(0.1)
    tf_tc.margin_bottom = Inches(0.1)
    
    # Category label
    p_cat = tf_tc.paragraphs[0]
    p_cat.text = tc["category"]
    p_cat.font.name = 'Manrope'
    p_cat.font.size = Pt(8)
    p_cat.font.bold = True
    p_cat.font.color.rgb = TEXT_SLATE
    p_cat.space_after = Pt(3)
    
    # Technology names
    p_tech = tf_tc.add_paragraph()
    p_tech.text = tc["techs"]
    p_tech.font.name = 'Manrope'
    p_tech.font.size = Pt(12)
    p_tech.font.bold = True
    p_tech.font.color.rgb = TEXT_CYAN
    p_tech.space_after = Pt(4)
    
    # Reason
    p_reason = tf_tc.add_paragraph()
    p_reason.text = tc["reason"]
    p_reason.font.name = 'Manrope'
    p_reason.font.size = Pt(8.5)
    p_reason.font.color.rgb = TEXT_WHITE

# 6. Footer compliance note
footer9_box = slide_9.shapes.add_textbox(Inches(0.41), Inches(4.80), Inches(9.32), Inches(0.2))
tf_footer9 = footer9_box.text_frame
tf_footer9.word_wrap = True
tf_footer9.margin_left = tf_footer9.margin_top = tf_footer9.margin_right = tf_footer9.margin_bottom = 0
p_footer9 = tf_footer9.paragraphs[0]
p_footer9.alignment = PP_ALIGN.CENTER
p_footer9.text = "Runs fully offline. No cloud services. No external APIs."
p_footer9.font.name = 'Manrope'
p_footer9.font.size = Pt(9)
p_footer9.font.bold = True
p_footer9.font.color.rgb = TEXT_CYAN

print(f"Saving final presentation to {OUTPUT_PPTX}...")
prs.save(OUTPUT_PPTX)
print("Presentation saved successfully.")

# ==============================================================================
# SCREENSHOT EXPORTING VIA POWERPOINT COM
# ==============================================================================
print("Launching PowerPoint to export slide screenshots...")
try:
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    abs_pptx_path = os.path.abspath(OUTPUT_PPTX)
    presentation = ppt_app.Presentations.Open(abs_pptx_path, WithWindow=True)
    
    # Export Slides 2-9
    for slide_num in range(2, 10):
        slide_img = os.path.join(ARTIFACTS_DIR, f"slide_{slide_num}_screenshot.png")
        print(f"Exporting Slide {slide_num} to {slide_img}...")
        presentation.Slides(slide_num).Export(slide_img, "PNG")
    
    presentation.Close()
    ppt_app.Quit()
    print("PowerPoint COM operations completed successfully.")
except Exception as e:
    print(f"PowerPoint Automation Error: {e}")
