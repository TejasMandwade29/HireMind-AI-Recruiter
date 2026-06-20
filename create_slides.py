import os
import shutil
import win32com.client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- PATHS SETUP ---
WORKSPACE_DIR = r"c:\Users\Dell\Downloads\[PUB] India_runs_data_and_ai_challenge\[PUB] India_runs_data_and_ai_challenge\India_runs_data_and_ai_challenge"
ARTIFACTS_DIR = r"C:\Users\Dell\.gemini\antigravity-ide\brain\45c24265-7a76-4c4c-a037-19d50c2804c1"
SCREENSHOT_SRC = os.path.join(ARTIFACTS_DIR, "dashboard_p2e_profile_1781862776024.png")
SCREENSHOT_DST = os.path.join(WORKSPACE_DIR, "profile_screenshot.png")
OUTPUT_PPTX = os.path.join(WORKSPACE_DIR, "hiremind_presentation.pptx")

print("Copying dashboard screenshot to workspace...")
if os.path.exists(SCREENSHOT_SRC):
    shutil.copy(SCREENSHOT_SRC, SCREENSHOT_DST)
    print("Screenshot copied successfully.")
else:
    print(f"Error: Screenshot source {SCREENSHOT_SRC} does not exist!")

# --- STYLE CONSTANTS ---
BG_COLOR = RGBColor(15, 23, 42)        # Deep Navy (#0F172A)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Light Navy (#1E293B)
TEXT_WHITE = RGBColor(255, 255, 255)   # White
TEXT_CYAN = RGBColor(0, 242, 254)     # Electric Cyan (#00F2FE)
TEXT_SLATE = RGBColor(148, 163, 184)   # Slate Gray (#94A3B8)
TEXT_CRIMSON = RGBColor(239, 68, 68)   # Crimson (#EF4444)
CARD_BORDER_COLOR = RGBColor(71, 85, 105) # Slate Border (#475569)

# --- PRESENTATION GENERATION ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide Layout 6 is blank slide in default python-pptx templates
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_header(slide, title_text):
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_CYAN

# ==============================================================================
# SLIDE 1: Cover Slide
# ==============================================================================
print("Creating Slide 1 (Cover)...")
slide_1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_1)

# Cover Header container
cover_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
tf = cover_box.text_frame
tf.word_wrap = True

# Slide Title
p_title = tf.paragraphs[0]
p_title.text = "HireMind"
p_title.font.name = 'Trebuchet MS'
p_title.font.size = Pt(54)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_CYAN
p_title.space_after = Pt(20)

# Slide Subtitle / Tagline
p_tag = tf.add_paragraph()
p_tag.text = "From 100,000 Candidates to the Right Hire"
p_tag.font.name = 'Arial'
p_tag.font.size = Pt(22)
p_tag.font.color.rgb = TEXT_WHITE
p_tag.space_after = Pt(40)

# Slide Meta
p_meta = tf.add_paragraph()
p_meta.text = "Problem Statement: Identify, validate, rank, and explain the most suitable candidates using AI-driven hiring intelligence.\nTeam Leader Name: Tejas Mandwade"
p_meta.font.name = 'Arial'
p_meta.font.size = Pt(14)
p_meta.font.color.rgb = TEXT_SLATE

# ==============================================================================
# SLIDE 2: Solution Overview
# ==============================================================================
print("Creating Slide 2 (Solution Overview)...")
slide_2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_2)
add_slide_header(slide_2, "Solution Overview")

# Hero Headline & Tagline Below Title
hero_box = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.8))
tf_hero = hero_box.text_frame
tf_hero.word_wrap = True
tf_hero.margin_left = tf_hero.margin_top = tf_hero.margin_right = tf_hero.margin_bottom = 0

p_hero = tf_hero.paragraphs[0]
p_hero.text = "From 100,000 Candidates to the Right Hire"
p_hero.font.name = 'Trebuchet MS'
p_hero.font.size = Pt(22)
p_hero.font.bold = True
p_hero.font.color.rgb = TEXT_WHITE

p_tag = tf_hero.add_paragraph()
p_tag.text = "HireMind transforms massive, noisy applicant pools into recruiter-ready shortlists in seconds."
p_tag.font.name = 'Arial'
p_tag.font.size = Pt(13)
p_tag.font.color.rgb = TEXT_SLATE
p_tag.space_before = Pt(4)

# Left Column: SaaS Contrast Container Shape
left_card = slide_2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(6.2), Inches(4.7)
)
left_card.fill.solid()
left_card.fill.fore_color.rgb = CARD_BG_COLOR
left_card.line.color.rgb = TEXT_CYAN
left_card.line.width = Pt(1.5)

# Text Box Over Container
left_text = slide_2.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.3))
tf_left = left_text.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0

p_card_title = tf_left.paragraphs[0]
p_card_title.text = "TRADITIONAL ATS VS. HIREMIND"
p_card_title.font.name = 'Arial'
p_card_title.font.size = Pt(14)
p_card_title.font.bold = True
p_card_title.font.color.rgb = TEXT_CYAN
p_card_title.space_after = Pt(14)

# Traditional ATS Bullets
p_ats_title = tf_left.add_paragraph()
p_ats_title.text = "●  Traditional ATS"
p_ats_title.font.name = 'Arial'
p_ats_title.font.size = Pt(13)
p_ats_title.font.bold = True
p_ats_title.font.color.rgb = TEXT_CRIMSON
p_ats_title.space_after = Pt(4)

p_ats_b1 = tf_left.add_paragraph()
p_ats_b1.text = "    - Relies on naive keyword matching."
p_ats_b1.font.name = 'Arial'
p_ats_b1.font.size = Pt(12)
p_ats_b1.font.color.rgb = TEXT_SLATE

p_ats_b2 = tf_left.add_paragraph()
p_ats_b2.text = "    - Easily gamed by keyword-stuffed resumes."
p_ats_b2.font.name = 'Arial'
p_ats_b2.font.size = Pt(12)
p_ats_b2.font.color.rgb = TEXT_SLATE

p_ats_b3 = tf_left.add_paragraph()
p_ats_b3.text = "    - Fails to verify profile integrity."
p_ats_b3.font.name = 'Arial'
p_ats_b3.font.size = Pt(12)
p_ats_b3.font.color.rgb = TEXT_SLATE
p_ats_b3.space_after = Pt(18)

# HireMind Bullets
p_hm_title = tf_left.add_paragraph()
p_hm_title.text = "●  HireMind Engine"
p_hm_title.font.name = 'Arial'
p_hm_title.font.size = Pt(13)
p_hm_title.font.bold = True
p_hm_title.font.color.rgb = TEXT_CYAN
p_hm_title.space_after = Pt(4)

p_hm_b1 = tf_left.add_paragraph()
p_hm_b1.text = "    - Matches resumes to real career context."
p_hm_b1.font.name = 'Arial'
p_hm_b1.font.size = Pt(12)
p_hm_b1.font.color.rgb = TEXT_WHITE

p_hm_b2 = tf_left.add_paragraph()
p_hm_b2.text = "    - Neutralizes keyword gaming."
p_hm_b2.font.name = 'Arial'
p_hm_b2.font.size = Pt(12)
p_hm_b2.font.color.rgb = TEXT_WHITE

p_hm_b3 = tf_left.add_paragraph()
p_hm_b3.text = "    - Automatically identifies anomalous profiles."
p_hm_b3.font.name = 'Arial'
p_hm_b3.font.size = Pt(12)
p_hm_b3.font.color.rgb = TEXT_WHITE

# Right Column: Distillation Funnel Timeline Nodes
funnel_steps = [
    {"label": "100,000 Raw Applicants", "desc": "Initial unstructured applicant pool", "color": CARD_BORDER_COLOR, "txt_color": TEXT_WHITE},
    {"label": "Automated Anomaly & Honeypot Detection", "desc": "Filters fraudulent and conflicting profiles", "color": TEXT_CRIMSON, "txt_color": TEXT_WHITE},
    {"label": "81,101 Verified Profiles", "desc": "Clean candidate pool ready for evaluation", "color": CARD_BORDER_COLOR, "txt_color": TEXT_WHITE},
    {"label": "Recruiter-Grade Ranking Engine", "desc": "Applies 5 multi-signal weights", "color": CARD_BG_COLOR, "txt_color": TEXT_WHITE},
    {"label": "Top 100 Shortlisted", "desc": "Verified high-signal hires visible on dashboard", "color": TEXT_CYAN, "txt_color": BG_COLOR}
]

for idx, step in enumerate(funnel_steps):
    y_pos = 2.1 + (idx * 0.95)
    
    # Node Background shape
    node = slide_2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(y_pos), Inches(5.5), Inches(0.65)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = step["color"]
    node.line.fill.background()
    
    # Text Box over Node
    node_text = slide_2.shapes.add_textbox(Inches(7.4), Inches(y_pos + 0.05), Inches(5.3), Inches(0.55))
    tf_node = node_text.text_frame
    tf_node.word_wrap = True
    tf_node.margin_left = tf_node.margin_top = tf_node.margin_right = tf_node.margin_bottom = 0
    
    p_lbl = tf_node.paragraphs[0]
    p_lbl.text = step["label"]
    p_lbl.font.name = 'Arial'
    p_lbl.font.size = Pt(11)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = step["txt_color"]
    
    p_desc = tf_node.add_paragraph()
    p_desc.text = step["desc"]
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(9.5)
    p_desc.font.color.rgb = step["txt_color"] if step["color"] != TEXT_CYAN else BG_COLOR
    p_desc.space_before = Pt(1)
    
    # Add downward arrow between boxes
    if idx < 4:
        arrow_box = slide_2.shapes.add_textbox(Inches(9.85), Inches(y_pos + 0.65), Inches(0.5), Inches(0.3))
        tf_arrow = arrow_box.text_frame
        tf_arrow.margin_left = tf_arrow.margin_top = tf_arrow.margin_right = tf_arrow.margin_bottom = 0
        p_arr = tf_arrow.paragraphs[0]
        p_arr.text = "▼"
        p_arr.font.size = Pt(9)
        p_arr.font.color.rgb = TEXT_SLATE

# ==============================================================================
# SLIDE 3: Beyond Keyword Matching
# ==============================================================================
print("Creating Slide 3 (Beyond Keyword Matching)...")
slide_3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_3)
add_slide_header(slide_3, "Beyond Keyword Matching")

# Hero Subheader
sub_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
tf_sub = sub_box.text_frame
tf_sub.word_wrap = True
tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0

p_sub = tf_sub.paragraphs[0]
p_sub.text = "Candidates are evaluated across five recruiter-grade signals."
p_sub.font.name = 'Arial'
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = TEXT_WHITE

# Left Column: UI Screenshot Frame & Caption
img_border = slide_3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(5.6), Inches(4.3)
)
img_border.fill.solid()
img_border.fill.fore_color.rgb = CARD_BG_COLOR
img_border.line.color.rgb = TEXT_CYAN
img_border.line.width = Pt(1.5)

# Insert Screenshot inside the frame
if os.path.exists(SCREENSHOT_DST):
    print("Inserting dashboard screenshot on Slide 3...")
    slide_3.shapes.add_picture(
        SCREENSHOT_DST, Inches(0.55), Inches(1.95), Inches(5.5), Inches(4.2)
    )
else:
    print("Warning: Profile screenshot not found, skipping picture insert.")

# Add Screenshot Caption Box
caption_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(5.6), Inches(0.4))
tf_cap = caption_box.text_frame
tf_cap.word_wrap = True
tf_cap.margin_left = tf_cap.margin_top = tf_cap.margin_right = tf_cap.margin_bottom = 0

p_cap = tf_cap.paragraphs[0]
p_cap.text = "Actual UI: Visual score contribution breakdown for Rank #1 Candidate."
p_cap.font.name = 'Arial'
p_cap.font.size = Pt(10.5)
p_cap.font.italic = True
p_cap.font.color.rgb = TEXT_SLATE

# Right Column: The 5 Signal Cards
signal_cards = [
    {
        "title": "Template Fit", "weight": "25%",
        "desc": "Ranks role compatibility against target profiles rather than simple keyword matches."
    },
    {
        "title": "Semantic Skills", "weight": "25%",
        "desc": "Maps core skill depth and skill-group affinity using local vectors."
    },
    {
        "title": "Candidate Intent", "weight": "20%",
        "desc": "Evaluates active responsiveness, availability indicators, and notice periods."
    },
    {
        "title": "Growth Trajectory", "weight": "15%",
        "desc": "Tracks historical job title upgrades and promotion history."
    },
    {
        "title": "Tenure Stability", "weight": "15%",
        "desc": "Analyzes employment durations and filters out profile job-hopping."
    }
]

for idx, card in enumerate(signal_cards):
    y_pos = 1.9 + (idx * 0.96)
    
    # Card Background container shape
    card_shape = slide_3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(y_pos), Inches(6.333), Inches(0.85)
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = CARD_BG_COLOR
    card_shape.line.color.rgb = CARD_BORDER_COLOR
    card_shape.line.width = Pt(1.0)
    
    # Text box inside Card
    card_text = slide_3.shapes.add_textbox(Inches(6.7), Inches(y_pos + 0.08), Inches(5.933), Inches(0.69))
    tf_card = card_text.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = tf_card.margin_top = tf_card.margin_right = tf_card.margin_bottom = 0
    
    p_header = tf_card.paragraphs[0]
    p_header.text = f"{card['title']}  "
    p_header.font.name = 'Arial'
    p_header.font.size = Pt(12)
    p_header.font.bold = True
    p_header.font.color.rgb = TEXT_WHITE
    
    # Add colored weight label to same header paragraph
    run_weight = p_header.add_run()
    run_weight.text = f"[{card['weight']}]"
    run_weight.font.name = 'Arial'
    run_weight.font.size = Pt(12)
    run_weight.font.bold = True
    run_weight.font.color.rgb = TEXT_CYAN
    
    p_desc = tf_card.add_paragraph()
    p_desc.text = card["desc"]
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(10.5)
    p_desc.font.color.rgb = TEXT_SLATE
    p_desc.space_before = Pt(2)

print(f"Saving final presentation to {OUTPUT_PPTX}...")
prs.save(OUTPUT_PPTX)
print("Presentation saved successfully.")

# ==============================================================================
# SCREENSHOT EXPORTING VIA POWERPOINT COM
# ==============================================================================
print("Launching PowerPoint to export slide screenshots...")
try:
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    # Open presentation (must be absolute path)
    abs_pptx_path = os.path.abspath(OUTPUT_PPTX)
    print(f"Opening presentation: {abs_pptx_path}")
    presentation = ppt_app.Presentations.Open(abs_pptx_path, WithWindow=True)
    
    # Export Slide 2 and Slide 3
    slide_2_img = os.path.join(ARTIFACTS_DIR, "slide_2_screenshot.png")
    slide_3_img = os.path.join(ARTIFACTS_DIR, "slide_3_screenshot.png")
    
    # Slide index in PowerPoint COM is 1-indexed
    print(f"Exporting Slide 2 to {slide_2_img}...")
    presentation.Slides(2).Export(slide_2_img, "PNG")
    
    print(f"Exporting Slide 3 to {slide_3_img}...")
    presentation.Slides(3).Export(slide_3_img, "PNG")
    
    presentation.Close()
    ppt_app.Quit()
    print("PowerPoint COM operations completed successfully.")
except Exception as e:
    print(f"PowerPoint Automation Error: {e}")
