import os
import shutil
import win32com.client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --- PATHS SETUP ---
WORKSPACE_DIR = r"c:\Users\Dell\Downloads\[PUB] India_runs_data_and_ai_challenge\[PUB] India_runs_data_and_ai_challenge\India_runs_data_and_ai_challenge"
ARTIFACTS_DIR = r"C:\Users\Dell\.gemini\antigravity-ide\brain\45c24265-7a76-4c4c-a037-19d50c2804c1"
SCREENSHOT_DST = os.path.join(WORKSPACE_DIR, "profile_screenshot.png")
OUTPUT_PPTX = os.path.join(WORKSPACE_DIR, "hiremind_presentation.pptx")

# --- STYLE CONSTANTS ---
BG_COLOR = RGBColor(15, 23, 42)        # Deep Navy (#0F172A)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Light Navy (#1E293B)
TEXT_WHITE = RGBColor(255, 255, 255)   # White
TEXT_CYAN = RGBColor(0, 242, 254)     # Electric Cyan (#00F2FE)
TEXT_SLATE = RGBColor(148, 163, 184)   # Slate Gray (#94A3B8)
TEXT_CRIMSON = RGBColor(239, 68, 68)   # Crimson (#EF4444)
CARD_BORDER_COLOR = RGBColor(71, 85, 105) # Slate Border (#475569)

# --- PRESENTATION SETUP ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_header(slide, title_text):
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_CYAN

def add_metric_badge(slide, badge_text):
    # Add a pill shape for consistent metric badge
    badge_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.45), Inches(3.033), Inches(0.45)
    )
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = CARD_BG_COLOR
    badge_bg.line.color.rgb = TEXT_CYAN
    badge_bg.line.width = Pt(1.0)
    
    tf = badge_bg.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = badge_text.upper()
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_CYAN

# ==============================================================================
# SLIDE 1: Cover Slide
# ==============================================================================
print("Creating Slide 1 (Cover)...")
slide_1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_1)

cover_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
tf = cover_box.text_frame
tf.word_wrap = True

p_title = tf.paragraphs[0]
p_title.text = "HireMind"
p_title.font.name = 'Trebuchet MS'
p_title.font.size = Pt(54)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_CYAN
p_title.space_after = Pt(20)

p_tag = tf.add_paragraph()
p_tag.text = "From 100,000 Candidates to the Right Hire"
p_tag.font.name = 'Arial'
p_tag.font.size = Pt(22)
p_tag.font.color.rgb = TEXT_WHITE
p_tag.space_after = Pt(40)

p_meta = tf.add_paragraph()
p_meta.text = "Problem Statement: Identify, validate, rank, and explain the most suitable candidates using AI-driven hiring intelligence.\nTeam Leader Name: Tejas Mandwade"
p_meta.font.name = 'Arial'
p_meta.font.size = Pt(14)
p_meta.font.color.rgb = TEXT_SLATE

# ==============================================================================
# SLIDE 2: Solution Overview
# ==============================================================================
print("Creating Slide 2 (Solution Overview)...")
slide_2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_2)
add_slide_header(slide_2, "Solution Overview")
add_metric_badge(slide_2, "100,000 Profiles Processed")

# Hero Headline & Tagline
hero_box = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.8))
tf_hero = hero_box.text_frame
tf_hero.word_wrap = True
tf_hero.margin_left = tf_hero.margin_top = tf_hero.margin_right = tf_hero.margin_bottom = 0

p_hero = tf_hero.paragraphs[0]
p_hero.text = "From 100,000 Candidates to the Right Hire"
p_hero.font.name = 'Trebuchet MS'
p_hero.font.size = Pt(22)
p_hero.font.bold = True
p_hero.font.color.rgb = TEXT_WHITE

p_tag = tf_hero.add_paragraph()
p_tag.text = "HireMind transforms massive, noisy applicant pools into recruiter-ready shortlists in seconds."
p_tag.font.name = 'Arial'
p_tag.font.size = Pt(13)
p_tag.font.color.rgb = TEXT_SLATE
p_tag.space_before = Pt(4)

# Left Column - Side-by-Side Comparison Columns
# Column A: Traditional ATS Card
ats_card = slide_2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(3.0), Inches(4.7)
)
ats_card.fill.solid()
ats_card.fill.fore_color.rgb = CARD_BG_COLOR
ats_card.line.color.rgb = TEXT_CRIMSON
ats_card.line.width = Pt(1.5)

tf_ats = ats_card.text_frame
tf_ats.word_wrap = True
tf_ats.margin_left = Inches(0.2)
tf_ats.margin_right = Inches(0.2)
tf_ats.margin_top = Inches(0.3)

p_ats_h = tf_ats.paragraphs[0]
p_ats_h.alignment = PP_ALIGN.CENTER
p_ats_h.text = "TRADITIONAL ATS"
p_ats_h.font.name = 'Arial'
p_ats_h.font.size = Pt(14)
p_ats_h.font.bold = True
p_ats_h.font.color.rgb = TEXT_CRIMSON
p_ats_h.space_after = Pt(24)

ats_rows = [
    ("Keyword Matching", "Simple keyword filtering"),
    ("Easy to Game", "Susceptible to text stuffing"),
    ("Fake Profiles Slip Through", "No profile anomaly filters"),
    ("Black Box Ranking", "Output results are unexplained")
]
for title, desc in ats_rows:
    p_tr = tf_ats.add_paragraph()
    p_tr.text = title
    p_tr.font.name = 'Arial'
    p_tr.font.size = Pt(12)
    p_tr.font.bold = True
    p_tr.font.color.rgb = TEXT_WHITE
    
    p_dr = tf_ats.add_paragraph()
    p_dr.text = desc
    p_dr.font.name = 'Arial'
    p_dr.font.size = Pt(10)
    p_dr.font.color.rgb = TEXT_SLATE
    p_dr.space_after = Pt(12)

# Column B: HireMind Card
hm_card = slide_2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(3.7), Inches(2.1), Inches(3.0), Inches(4.7)
)
hm_card.fill.solid()
hm_card.fill.fore_color.rgb = CARD_BG_COLOR
hm_card.line.color.rgb = TEXT_CYAN
hm_card.line.width = Pt(1.5)

tf_hm = hm_card.text_frame
tf_hm.word_wrap = True
tf_hm.margin_left = Inches(0.2)
tf_hm.margin_right = Inches(0.2)
tf_hm.margin_top = Inches(0.3)

p_hm_h = tf_hm.paragraphs[0]
p_hm_h.alignment = PP_ALIGN.CENTER
p_hm_h.text = "HIREMIND"
p_hm_h.font.name = 'Arial'
p_hm_h.font.size = Pt(14)
p_hm_h.font.bold = True
p_hm_h.font.color.rgb = TEXT_CYAN
p_hm_h.space_after = Pt(24)

hm_rows = [
    ("Multi-Signal Ranking", "Evaluates five signals"),
    ("Detects Manipulation", "Cross-validates history"),
    ("Automated Anomaly Detection", "Filters honeypots & spam"),
    ("Explainable Decisions", "Visual, transparent reasoning")
]
for title, desc in hm_rows:
    p_tr = tf_hm.add_paragraph()
    p_tr.text = title
    p_tr.font.name = 'Arial'
    p_tr.font.size = Pt(12)
    p_tr.font.bold = True
    p_tr.font.color.rgb = TEXT_CYAN
    
    p_dr = tf_hm.add_paragraph()
    p_dr.text = desc
    p_dr.font.name = 'Arial'
    p_dr.font.size = Pt(10)
    p_dr.font.color.rgb = TEXT_WHITE
    p_dr.space_after = Pt(12)

# Right Column: The Funnel with HERO Numbers
funnel_steps = [
    {"num": "100,000", "label": "RAW APPLICANTS", "color": CARD_BORDER_COLOR, "num_color": TEXT_WHITE, "txt_color": TEXT_SLATE, "height": 0.65},
    {"num": "81,101", "label": "CLEAN CANDIDATES (Anomaly-free)", "color": CARD_BORDER_COLOR, "num_color": TEXT_WHITE, "txt_color": TEXT_SLATE, "height": 0.65},
    {"num": "1,092", "label": "TIER 1 & 2 SPECIALISTS MATCHED", "color": CARD_BG_COLOR, "num_color": TEXT_WHITE, "txt_color": TEXT_SLATE, "height": 0.65},
    {"num": "🏆 TOP 100", "label": "RECRUITER-READY SHORTLIST", "color": TEXT_CYAN, "num_color": BG_COLOR, "txt_color": BG_COLOR, "height": 1.3}
]

y_pos = 2.1
for idx, step in enumerate(funnel_steps):
    box_h = step["height"]
    
    # Node Shape
    node = slide_2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(y_pos), Inches(5.5), Inches(box_h)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = step["color"]
    node.line.fill.background()
    
    if step["color"] == TEXT_CYAN:
        node.line.color.rgb = TEXT_CYAN
        node.line.width = Pt(2.0)
        
    # Text Box over Node
    node_text = slide_2.shapes.add_textbox(Inches(7.4), Inches(y_pos + (0.04 if box_h < 1.0 else 0.2)), Inches(5.3), Inches(box_h - 0.1))
    tf_node = node_text.text_frame
    tf_node.word_wrap = True
    tf_node.margin_left = tf_node.margin_top = tf_node.margin_right = tf_node.margin_bottom = 0
    
    p_num = tf_node.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = step["num"]
    p_num.font.name = 'Arial'
    p_num.font.size = Pt(26) if box_h < 1.0 else Pt(32)
    p_num.font.bold = True
    p_num.font.color.rgb = step["num_color"]
    
    p_lbl = tf_node.add_paragraph()
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.text = step["label"]
    p_lbl.font.name = 'Arial'
    p_lbl.font.size = Pt(9.5) if box_h < 1.0 else Pt(11)
    p_lbl.font.bold = True if box_h >= 1.0 else False
    p_lbl.font.color.rgb = step["txt_color"]
    p_lbl.space_before = Pt(1)
    
    # Downward arrow between boxes
    if idx < 3:
        arrow_box = slide_2.shapes.add_textbox(Inches(9.85), Inches(y_pos + box_h), Inches(0.5), Inches(0.3))
        tf_arrow = arrow_box.text_frame
        tf_arrow.margin_left = tf_arrow.margin_top = tf_arrow.margin_right = tf_arrow.margin_bottom = 0
        p_arr = tf_arrow.paragraphs[0]
        p_arr.alignment = PP_ALIGN.CENTER
        p_arr.text = "▼"
        p_arr.font.size = Pt(9)
        p_arr.font.color.rgb = TEXT_SLATE
        
    y_pos += box_h + 0.35

# ==============================================================================
# SLIDE 3: Beyond Keyword Matching
# ==============================================================================
print("Creating Slide 3 (Beyond Keyword Matching)...")
slide_3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide_3)
add_slide_header(slide_3, "Beyond Keyword Matching")
add_metric_badge(slide_3, "5 Recruiter Signals Evaluated")

# Hero Subheader
sub_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
tf_sub = sub_box.text_frame
tf_sub.word_wrap = True
tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0

p_sub = tf_sub.paragraphs[0]
p_sub.text = "Candidates are evaluated across five recruiter-grade signals."
p_sub.font.name = 'Arial'
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = TEXT_WHITE

# Left Column - High Resolution Screenshot Crop Frame (Widescreen layout)
img_border = slide_3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.85), Inches(5.8), Inches(2.5)
)
img_border.fill.solid()
img_border.fill.fore_color.rgb = CARD_BG_COLOR
img_border.line.color.rgb = TEXT_CYAN
img_border.line.width = Pt(1.5)

# Insert Screenshot inside the frame
if os.path.exists(SCREENSHOT_DST):
    print("Inserting cropped screenshot on Slide 3...")
    slide_3.shapes.add_picture(
        SCREENSHOT_DST, Inches(0.52), Inches(1.87), Inches(5.76), Inches(2.46)
    )
else:
    print("Warning: profile_screenshot.png not found, skipping image insert.")

caption_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(4.45), Inches(5.8), Inches(0.3))
tf_cap = caption_box.text_frame
tf_cap.word_wrap = True
tf_cap.margin_left = tf_cap.margin_top = tf_cap.margin_right = tf_cap.margin_bottom = 0
p_cap = tf_cap.paragraphs[0]
p_cap.text = "Actual UI: Score contribution visualizer for Rank #1 Candidate."
p_cap.font.name = 'Arial'
p_cap.font.size = Pt(10.5)
p_cap.font.italic = True
p_cap.font.color.rgb = TEXT_SLATE

# Left Column Bottom: Trust Engine Capabilities Callout Card
trust_card = slide_3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.85), Inches(5.8), Inches(1.95)
)
trust_card.fill.solid()
trust_card.fill.fore_color.rgb = CARD_BG_COLOR
trust_card.line.color.rgb = TEXT_CYAN
trust_card.line.width = Pt(1.0)

tf_tr = trust_card.text_frame
tf_tr.word_wrap = True
tf_tr.margin_left = Inches(0.2)
tf_tr.margin_top = Inches(0.18)

p_tr_h = tf_tr.paragraphs[0]
p_tr_h.text = "TRUST ENGINE CAPABILITIES"
p_tr_h.font.name = 'Arial'
p_tr_h.font.size = Pt(12)
p_tr_h.font.bold = True
p_tr_h.font.color.rgb = TEXT_CYAN
p_tr_h.space_after = Pt(12)

p_tr_l1 = tf_tr.add_paragraph()
p_tr_l1.text = "✓  Skill Depth    |    ✓  Career Growth    |    ✓  Role Fit"
p_tr_l1.font.name = 'Arial'
p_tr_l1.font.size = Pt(11)
p_tr_l1.font.bold = True
p_tr_l1.font.color.rgb = TEXT_WHITE
p_tr_l1.space_after = Pt(8)

p_tr_l2 = tf_tr.add_paragraph()
p_tr_l2.text = "✓  Recruiter Engagement    |    ✓  Profile Stability"
p_tr_l2.font.name = 'Arial'
p_tr_l2.font.size = Pt(11)
p_tr_l2.font.bold = True
p_tr_l2.font.color.rgb = TEXT_WHITE

p_tr_footer = tf_tr.add_paragraph()
p_tr_footer.text = "Engine cross-references stated skills with response rate, active days and job tenures."
p_tr_footer.font.name = 'Arial'
p_tr_footer.font.size = Pt(9.5)
p_tr_footer.font.color.rgb = TEXT_SLATE
p_tr_footer.space_before = Pt(12)

# Right Column: The 5 Signal Cards (Vertical Layout)
signal_cards = [
    {
        "title": "Template Fit", "weight": "25%",
        "desc": "Role Compatibility (target profile match)"
    },
    {
        "title": "Semantic Skills", "weight": "25%",
        "desc": "Skill Alignment (embedding vector match)"
    },
    {
        "title": "Candidate Intent", "weight": "20%",
        "desc": "Recruiter Engagement (response speed & active flags)"
    },
    {
        "title": "Growth", "weight": "15%",
        "desc": "Career Progression (upgrades in title and tenure)"
    },
    {
        "title": "Stability", "weight": "15%",
        "desc": "Profile Stability (tenure duration check)"
    }
]

for idx, card in enumerate(signal_cards):
    y_pos = 1.85 + (idx * 0.98)
    
    # Card Background container shape
    card_shape = slide_3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(y_pos), Inches(6.333), Inches(0.85)
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = CARD_BG_COLOR
    card_shape.line.color.rgb = CARD_BORDER_COLOR
    card_shape.line.width = Pt(1.0)
    
    # Text box inside Card
    card_text = slide_3.shapes.add_textbox(Inches(6.7), Inches(y_pos + 0.12), Inches(5.933), Inches(0.65))
    tf_card = card_text.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = tf_card.margin_top = tf_card.margin_right = tf_card.margin_bottom = 0
    
    p_header = tf_card.paragraphs[0]
    
    # Add colored weight label as primary visual element
    run_weight = p_header.add_run()
    run_weight.text = f"{card['weight']}   "
    run_weight.font.name = 'Arial'
    run_weight.font.size = Pt(20)
    run_weight.font.bold = True
    run_weight.font.color.rgb = TEXT_CYAN
    
    # Add slide title next to it
    run_title = p_header.add_run()
    run_title.text = card['title']
    run_title.font.name = 'Arial'
    run_title.font.size = Pt(13)
    run_title.font.bold = True
    run_title.font.color.rgb = TEXT_WHITE
    
    p_desc = tf_card.add_paragraph()
    p_desc.text = card["desc"]
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(10.5)
    p_desc.font.color.rgb = TEXT_SLATE
    p_desc.space_before = Pt(1)

print(f"Saving presentation to {OUTPUT_PPTX}...")
prs.save(OUTPUT_PPTX)
print("Presentation saved successfully.")

# ==============================================================================
# SCREENSHOT EXPORTING VIA POWERPOINT COM
# ==============================================================================
print("Launching PowerPoint to export slide screenshots...")
try:
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    abs_pptx_path = os.path.abspath(OUTPUT_PPTX)
    presentation = ppt_app.Presentations.Open(abs_pptx_path, WithWindow=True)
    
    slide_2_img = os.path.join(ARTIFACTS_DIR, "slide_2_screenshot.png")
    slide_3_img = os.path.join(ARTIFACTS_DIR, "slide_3_screenshot.png")
    
    print(f"Exporting Slide 2 to {slide_2_img}...")
    presentation.Slides(2).Export(slide_2_img, "PNG")
    
    print(f"Exporting Slide 3 to {slide_3_img}...")
    presentation.Slides(3).Export(slide_3_img, "PNG")
    
    presentation.Close()
    ppt_app.Quit()
    print("PowerPoint COM operations completed successfully.")
except Exception as e:
    print(f"PowerPoint Automation Error: {e}")
